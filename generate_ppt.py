import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen Layout (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette (Neo-Brutalism Theme)
    C_DARK = RGBColor(18, 18, 18)       # #121212
    C_YELLOW = RGBColor(255, 230, 0)    # #ffe600
    C_GREEN = RGBColor(167, 243, 208)   # #a7f3d0
    C_BLUE = RGBColor(112, 214, 255)    # #70d6ff
    C_PINK = RGBColor(255, 153, 200)    # #ff99c8
    C_BG = RGBColor(246, 245, 240)      # #f6f5f0
    C_WHITE = RGBColor(255, 255, 255)
    C_GRAY = RGBColor(100, 100, 100)

    blank_layout = prs.slide_layouts[6]

    def add_neo_card(slide, left, top, width, height, bg_color, border_color=C_DARK):
        # Shadow Shape
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.08), top + Inches(0.08), width, height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = C_DARK
        shadow.line.fill.background()

        # Main Shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(2.5)
        return card

    # ==================== SLIDE 1: COVER ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_BG
    bg1.line.fill.background()

    # Main Card Cover
    card1 = add_neo_card(slide1, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1), C_YELLOW)
    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "SISTEM ASPIRASI & PENGADUAN WARGA DESA"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = C_DARK

    p2 = tf1.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "LAPOR DESA RAU ★"
    r2.font.size = Pt(40)
    r2.font.bold = True
    r2.font.color.rgb = C_DARK

    p3 = tf1.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Aplikasi Pelaporan Cepat Berbasis GPS Geofencing (Radius 2.0 KM)\nKecamatan Kedung, Kabupaten Jepara"
    r3.font.size = Pt(16)
    r3.font.bold = True
    r3.font.color.rgb = C_DARK

    p4 = tf1.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "\nWeb: https://lapor-desa-rau.vercel.app  •  Biaya Server: Rp 0,- (Gratis Selamanya)"
    r4.font.size = Pt(12)
    r4.font.bold = True
    r4.font.color.rgb = C_DARK

    # ==================== SLIDE 2: MASALAH & SOLUSI ====================
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_BG
    bg2.line.fill.background()

    # Header Title
    head2 = add_neo_card(slide2, Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.9), C_DARK)
    tf_h2 = head2.text_frame
    tf_h2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph2 = tf_h2.paragraphs[0]
    rh2 = ph2.add_run()
    rh2.text = "KENAPA DESA RAU BUTUH APLIKASI INI?"
    rh2.font.size = Pt(20)
    rh2.font.bold = True
    rh2.font.color.rgb = C_YELLOW

    # Card 1: Masalah Lama
    c_prob = add_neo_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), C_PINK)
    tf_prob = c_prob.text_frame
    tf_prob.word_wrap = True
    p_p1 = tf_prob.paragraphs[0]
    r_p1 = p_p1.add_run()
    r_p1.text = "❌ MASALAH CARA LAMA:"
    r_p1.font.size = Pt(16)
    r_p1.font.bold = True
    r_p1.font.color.rgb = C_DARK

    items_prob = [
        "1. Warga sungkan / enggan lapor langsung ke kantor desa.",
        "2. Laporan sering tertimbun di chat WA grup warga dan terlupakan.",
        "3. Lokasi jalan/lampu rusak tidak akurat (petugas bingung mencari titik masalah).",
        "4. Tidak ada kejelasan apakah laporan sudah dikerjakan atau belum.",
        "5. Biaya pengadaan aplikasi desa biasanya sangat mahal."
    ]
    for ip in items_prob:
        pp = tf_prob.add_paragraph()
        rp = pp.add_run()
        rp.text = ip
        rp.font.size = Pt(12)
        rp.font.bold = True
        rp.font.color.rgb = C_DARK

    # Card 2: Solusi Lapor Rau
    c_sol = add_neo_card(slide2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), C_GREEN)
    tf_sol = c_sol.text_frame
    tf_sol.word_wrap = True
    p_s1 = tf_sol.paragraphs[0]
    r_s1 = p_s1.add_run()
    r_s1.text = "✅ SOLUSI LAPOR.RAU:"
    r_s1.font.size = Pt(16)
    r_s1.font.bold = True
    r_s1.font.color.rgb = C_DARK

    items_sol = [
        "1. Lapor 24 jam dari HP tanpa repot daftar akun/password.",
        "2. Setiap aduan punya Kode Tiket resmi & masuk antrean sistem.",
        "3. GPS HP otomatis mendeteksi titik koordinat presisi di peta.",
        "4. Transparansi publik: Warga bisa memantau status sampai 'BERES'.",
        "5. Biaya server Rp 0,- (Zero-Cost Cloud Architecture)."
    ]
    for isol in items_sol:
        ps = tf_sol.add_paragraph()
        rs = ps.add_run()
        rs.text = isol
        rs.font.size = Pt(12)
        rs.font.bold = True
        rs.font.color.rgb = C_DARK

    # ==================== SLIDE 3: 5 FITUR UNGGULAN ====================
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = C_BG
    bg3.line.fill.background()

    head3 = add_neo_card(slide3, Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.9), C_YELLOW)
    tf_h3 = head3.text_frame
    tf_h3.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph3 = tf_h3.paragraphs[0]
    rh3 = ph3.add_run()
    rh3.text = "FITUR-FITUR UNGGULAN UNTUK WARGA"
    rh3.font.size = Pt(20)
    rh3.font.bold = True
    rh3.font.color.rgb = C_DARK

    # 4 Bento Cards
    bento_items = [
        ("📍 RADAR GPS GEOFENCE", "Hanya warga dalam radius 2.0 km dari Balai Desa Rau yang bisa kirim laporan. Anti-hoaks dari orang luar.", C_BLUE, Inches(0.8), Inches(1.8)),
        ("⚡ TANPA LOGIN & RAMAH LANSIA", "Cukup buka website di browser HP, langsung siap pakai. Bisa dipasang ke layar utama HP (PWA).", C_WHITE, Inches(6.9), Inches(1.8)),
        ("📷 KOMPRESI FOTO OTOMATIS", "Foto kamera HP 8 MB otomatis diperkecil jadi format WebP ~50 KB. Sangat hemat kuota dan kirim super cepat.", C_WHITE, Inches(0.8), Inches(4.3)),
        ("🕵️ OPSI KIRIM ANONIM", "Warga yang sungkan namanya terlihat tetangga bisa memilih mode anonim. Identitas terlindungi aman.", C_PINK, Inches(6.9), Inches(4.3)),
    ]

    for title, desc, col, l, t in bento_items:
        card_b = add_neo_card(slide3, l, t, Inches(5.6), Inches(2.2), col)
        tf_b = card_b.text_frame
        tf_b.word_wrap = True
        pb1 = tf_b.paragraphs[0]
        rb1 = pb1.add_run()
        rb1.text = title
        rb1.font.size = Pt(14)
        rb1.font.bold = True
        rb1.font.color.rgb = C_DARK

        pb2 = tf_b.add_paragraph()
        rb2 = pb2.add_run()
        rb2.text = desc
        rb2.font.size = Pt(11)
        rb2.font.color.rgb = C_DARK

    # ==================== SLIDE 4: ALUR KERJA PETUGAS ====================
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = C_BG
    bg4.line.fill.background()

    head4 = add_neo_card(slide4, Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.9), C_DARK)
    tf_h4 = head4.text_frame
    tf_h4.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph4 = tf_h4.paragraphs[0]
    rh4 = ph4.add_run()
    rh4.text = "PORTAL KHUSUS PETUGAS BALAI DESA (/kantor-desa)"
    rh4.font.size = Pt(19)
    rh4.font.bold = True
    rh4.font.color.rgb = C_GREEN

    # Petugas Feature List
    c_petugas = add_neo_card(slide4, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8), C_WHITE)
    tf_p = c_petugas.text_frame
    tf_p.word_wrap = True

    p_p_title = tf_p.paragraphs[0]
    r_pt = p_p_title.add_run()
    r_pt.text = "Fitur Rahasia & Tersembunyi Khusus Staf Balai Desa (PIN: rau2026):"
    r_pt.font.size = Pt(15)
    r_pt.font.bold = True
    r_pt.font.color.rgb = C_DARK

    petugas_features = [
        "1. Triage Aduan: Mengubah status dari Menunggu ➔ Diproses ➔ Beres ✓ secara instan.",
        "2. Tanggapan Resmi: Mengisi pesan perkembangan penanganan yang langsung terbaca oleh warga.",
        "3. Upload Foto Bukti Perbaikan: Mengunggah foto hasil pengerjaan jalan/lampu yang sudah diperbaiki.",
        "4. Tombol Chat WA Otomatis: Satu klik langsung terhubung ke WhatsApp warga pelapor dengan template pesan santun bahasa Jawa krama / Indonesia.",
        "5. Ekspor Rekapitulasi CSV: Mengunduh data rekap bulanan ke format Excel/CSV untuk arsip laporan ke Kecamatan / BPD."
    ]

    for pf in petugas_features:
        ppf = tf_p.add_paragraph()
        rpf = ppf.add_run()
        rpf.text = pf
        rpf.font.size = Pt(12)
        rpf.font.bold = True
        rpf.font.color.rgb = C_DARK

    # ==================== SLIDE 5: CARA WARGA MELAPOR ====================
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = C_BG
    bg5.line.fill.background()

    head5 = add_neo_card(slide5, Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.9), C_YELLOW)
    tf_h5 = head5.text_frame
    tf_h5.vertical_anchor = MSO_ANCHOR.MIDDLE
    ph5 = tf_h5.paragraphs[0]
    rh5 = ph5.add_run()
    rh5.text = "PANDUAN 3 LANGKAH MUDAH UNTUK WARGA"
    rh5.font.size = Pt(20)
    rh5.font.bold = True
    rh5.font.color.rgb = C_DARK

    step_cards = [
        ("LANGKAH 1 📸", "Jepret Foto & Lokasi", "Buka https://lapor-desa-rau.vercel.app dari HP. Foto masalah di lokasi kejadian (GPS otomatis mencatat titik koordinat).", C_PINK, Inches(0.8)),
        ("LANGKAH 2 📝", "Isi Keterangan", "Pilih jenis masalah (Jalan, Lampu, Sampah, Bansos), pilih nomor RT/RW, dan klik tombol 'Kirim ke Balai Desa'.", C_BLUE, Inches(4.8)),
        ("LANGKAH 3 🔍", "Pantau Sampai Beres", "Simpan nomor tiket unik laporan Anda atau lihat di menu 'PANTAU'. Pantau sampai statusnya berubah jadi BERES ✓.", C_GREEN, Inches(8.8)),
    ]

    for stitle, ssub, sdesc, scol, sleft in step_cards:
        sc = add_neo_card(slide5, sleft, Inches(1.8), Inches(3.7), Inches(4.8), scol)
        tfc = sc.text_frame
        tfc.word_wrap = True

        p1 = tfc.paragraphs[0]
        r1 = p1.add_run()
        r1.text = stitle
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = C_DARK

        p2 = tfc.add_paragraph()
        r2 = p2.add_run()
        r2.text = ssub
        r2.font.size = Pt(13)
        r2.font.bold = True
        r2.font.color.rgb = C_DARK

        p3 = tfc.add_paragraph()
        r3 = p3.add_run()
        r3.text = "\n" + sdesc
        r3.font.size = Pt(11)
        r3.font.bold = True
        r3.font.color.rgb = C_DARK

    # ==================== SLIDE 6: PENUTUP ====================
    slide6 = prs.slides.add_slide(blank_layout)
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = C_BG
    bg6.line.fill.background()

    card6 = add_neo_card(slide6, Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.5), C_DARK)
    tf6 = card6.text_frame
    tf6.word_wrap = True
    tf6.vertical_anchor = MSO_ANCHOR.MIDDLE

    p6_1 = tf6.paragraphs[0]
    p6_1.alignment = PP_ALIGN.CENTER
    r6_1 = p6_1.add_run()
    r6_1.text = "MATUR NUWUN SANGET"
    r6_1.font.size = Pt(36)
    r6_1.font.bold = True
    r6_1.font.color.rgb = C_YELLOW

    p6_2 = tf6.add_paragraph()
    p6_2.alignment = PP_ALIGN.CENTER
    r6_2 = p6_2.add_run()
    r6_2.text = "\n\"Sareng-sareng Mbangun Desa Rau Supados Langkung Sae lan Maju\"\n"
    r6_2.font.size = Pt(16)
    r6_2.font.bold = True
    r6_2.font.color.rgb = C_WHITE

    p6_3 = tf6.add_paragraph()
    p6_3.alignment = PP_ALIGN.CENTER
    r6_3 = p6_3.add_run()
    r6_3.text = "Mari sosialisasikan ke grup WhatsApp RT dan warga Desa Rau:\n👉 https://lapor-desa-rau.vercel.app"
    r6_3.font.size = Pt(14)
    r6_3.font.bold = True
    r6_3.font.color.rgb = C_GREEN

    out_ppt = "D:/Lapor Desa Rau/Presentasi_Lapor_Desa_Rau.pptx"
    prs.save(out_ppt)
    print("SUCCESS: PowerPoint Presentation berhasil dibuat di " + out_ppt)

create_presentation()
